"""Загрузчики документов: каждый возвращает текст (+ при наличии — постранично).

Поддержка: PDF, DOCX, PPTX, XLSX/CSV (прайс-листы), TXT/MD, HTML,
аудио/видео (обучающие записи -> транскрибация Whisper).
"""
from __future__ import annotations
import warnings
from pathlib import Path
from typing import Iterator
import settings

# шумные предупреждения парсеров (openpyxl Data Validation и т.п.)
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
warnings.filterwarnings("ignore", message=".*Data Validation extension.*")

# Какие расширения к какому обработчику
AUDIO_VIDEO = {".mp3", ".wav", ".m4a", ".aac", ".mp4", ".mov", ".mkv", ".webm"}
# RAW-фото камер: конвертируются в изображение → OCR → текстовый PDF
RAW_PHOTO = {".cr2", ".cr3", ".nef", ".arw", ".dng", ".raf", ".rw2", ".orf", ".sr2"}
# Растровые изображения: распознавание текста (OCR)
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".jfif"}
# Архивы: распаковываются, содержимое индексируется как обычные файлы
ARCHIVE_EXTS = {".zip", ".rar", ".7z", ".tar", ".gz", ".tgz", ".bz2"}
_ARCHIVE_MAX_DEPTH = 2               # защита от вложенных архивов / «архивных бомб»
_ARCHIVE_MAX_FILES = 5000            # лимит файлов внутри одного архива
_ARCHIVE_MAX_BYTES = 4 * 1024 ** 3   # лимит суммарного распакованного объёма (4 ГБ)


def probe_file(path_str: str, timeout: int = 0):
    """Проверить, извлекается ли текст из файла. Возвращает (status, issue):
    status ∈ {'ok','failed','timeout'}, issue — текст проблемы или None.
    Используется пулом процессов для ПАРАЛЛЕЛЬНОЙ проверки каталога."""
    import signal
    p = Path(path_str)
    use_alarm = bool(timeout) and hasattr(signal, "SIGALRM")

    def _to(_s, _f):
        raise TimeoutError()
    try:
        if use_alarm:
            signal.signal(signal.SIGALRM, _to)
            signal.alarm(int(timeout))
        got = False
        for part in load_file(p):
            if (part.get("text", "") or "").strip():
                got = True
                break  # текст найден — дальше можно не читать
        if use_alarm:
            signal.alarm(0)
        return ("ok", None) if got else ("failed", "текст не извлечён")
    except TimeoutError:
        if use_alarm:
            signal.alarm(0)
        return ("timeout", f"превышен лимит {timeout} c")
    except Exception as e:
        if use_alarm:
            signal.alarm(0)
        return ("failed", str(e)[:200])


def load_file(path: Path, _depth: int = 0) -> Iterator[dict]:
    """Yield {'text', 'page'} для одного файла. Пустые куски пропускаются.
    `_depth` — внутренний счётчик вложенности для распаковки архивов."""
    ext = path.suffix.lower()
    try:
        if ext in ARCHIVE_EXTS:
            yield from _load_archive(path, _depth)
            return
        if ext == ".pdf":
            yield from _load_pdf(path)
        elif ext == ".docx":
            yield from _load_docx(path)
        elif ext == ".pptx":
            yield from _load_pptx(path)
        elif ext in {".xlsx", ".xlsm", ".xls", ".csv"}:
            yield from _load_table(path)
        elif ext in {".txt", ".md"}:
            yield {"text": _read_text_any(path), "page": None}
        elif ext in {".html", ".htm", ".mhtml", ".mht"}:
            yield from _load_html(path)
        elif ext == ".doc":
            yield from _load_doc(path)
        elif ext in {".xml"}:
            yield from _load_xml(path)
        elif ext == ".json":
            yield from _load_json(path)
        elif ext == ".url":
            yield from _load_url(path)
        elif ext == ".msg":
            yield from _load_msg(path)
        elif ext == ".svg":
            yield from _load_svg(path)
        elif ext in {".dxf", ".dwg", ".stp", ".step", ".igs", ".iges"}:
            # чертежи/3D-CAD — тяжёлая конвертация DWG; можно отключить ради скорости
            if _enabled("PARSE_CAD"):
                if ext in {".dxf", ".dwg"}:
                    yield from _load_cad(path)
                else:
                    yield from _load_cad_exchange(path)
        elif ext in IMAGE_EXTS:
            if _enabled("OCR_IMAGES"):   # OCR изображений — самый долгий этап
                yield from _load_image(path)
        elif ext in RAW_PHOTO:
            if _enabled("OCR_RAW"):
                yield from _load_raw(path)
        elif ext in AUDIO_VIDEO:
            if _enabled("TRANSCRIBE_AV"):  # транскрибация Whisper — минуты на файл
                yield from _load_av(path)
        # остальное молча пропускаем
    except Exception as e:
        # Ошибку парсера прокидываем наружу — в ingest она попадёт в список errors и
        # станет видна пользователю (раньше глоталась print-ом и файл «молча пустел»).
        # Подавляем только для вложенных файлов архива (_depth > 0): один битый файл
        # внутри архива не должен срывать индексацию остального содержимого.
        if _depth > 0:
            print(f"  ! ошибка чтения {path.name}: {e}")
            return
        raise


def _enabled(key: str) -> bool:
    """Включён ли тяжёлый экстрактор (по рантайм-настройке, по умолчанию True)."""
    try:
        v = settings.get(key)
        return True if v is None else bool(v)
    except Exception:
        return True


def _describe_due(total_chars: int) -> bool:
    """Нужно ли описывать изображение моделью: получилось ≤ N чанков (0/1 по умолч.)."""
    if not _enabled("OCR_LLM_DESCRIBE") or not _vision_available():
        return False
    try:
        cs = int(settings.get("CHUNK_SIZE") or 900)
    except Exception:
        cs = 900
    try:
        maxc = int(_ocr_setting("OCR_LLM_MAX_CHUNKS", 1))
    except Exception:
        maxc = 1
    import math
    chunks = 0 if total_chars <= 0 else math.ceil(total_chars / max(1, cs))
    return chunks <= maxc


# Предохранитель vision: если модель падает подряд (vLLM «Server disconnected»/краш движка),
# перестаём её дёргать до конца процесса — иначе тормозим индексацию и добиваем упавший vLLM.
_vision_state = {"fails": 0, "off": False}


def _vision_available() -> bool:
    return not _vision_state["off"]


def _vision_result(ok: bool) -> None:
    if ok:
        _vision_state["fails"] = 0
    else:
        _vision_state["fails"] += 1
        if _vision_state["fails"] >= 8 and not _vision_state["off"]:
            _vision_state["off"] = True
            print("  ! vision-модель не отвечает 8 раз подряд — отключаю описание "
                  "картинок/чертежей до конца индексации. Перезапустите vLLM "
                  "(docker restart rag_vllm) и проверьте память/лимит одновременных запросов.",
                  flush=True)


def _describe_image_part(img, page=None):
    """Описать изображение vision-моделью → {'text','page'} или None."""
    if not _vision_available():
        return None
    try:
        import llm_backend
        desc = llm_backend.describe_image(img)
        _vision_result(bool(desc and desc.strip()))
        if desc and desc.strip():
            return {"text": "Описание изображения (vision-модель):\n" + desc.strip(),
                    "page": page, "vision_desc": True}
    except Exception as e:
        _vision_result(False)
        print(f"  ~ описание изображения моделью не удалось: {e}")
    return None


def describe_file_llm(source: str, text: str) -> str:
    """Краткое LLM-описание файла по его тексту (для опции INDEX_LLM_DESCRIBE)."""
    text = (text or "").strip()
    if not text:
        return ""
    try:
        maxc = int(settings.get("INDEX_LLM_DESCRIBE_MAXCHARS") or 6000)
    except Exception:
        maxc = 6000
    snippet = text[:max(200, maxc)]
    prompt = (
        f"Ниже содержимое документа «{source}». Составь краткое деловое описание: о чём "
        f"документ, ключевые сущности (продукты, артикулы, параметры, стороны, даты, цены), "
        f"его назначение. 3–6 предложений, по-русски, без вступлений.\n\n"
        f"ДОКУМЕНТ:\n{snippet}")
    try:
        import llm_backend
        out = llm_backend.chat([{"role": "user", "content": prompt}],
                               temperature=0.2, kind="doc-describe", label=source)
        return (out or "").strip()
    except Exception as e:
        print(f"  ~ LLM-описание файла не удалось ({source}): {e}")
        return ""


def _ocr_pdf_page(page, i):
    """Распознать «картиночную» страницу PDF (текст нарисован графикой): рендерим
    страницу в изображение и прогоняем через OCR. Возвращает {'text','page'}.
    Если OCR дал мало текста — дополнительно описываем страницу vision-моделью."""
    try:
        import fitz  # pymupdf
        from PIL import Image
        try:
            scale = float(_ocr_setting("OCR_SCALE", 2.5))
        except Exception:
            scale = 2.5
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale))  # масштаб из настроек
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        total = 0
        for part in _ocr_image(img):
            t = (part.get("text") or "").strip()
            if t:
                total += len(t)
                yield {"text": part["text"], "page": i}
        if _describe_due(total):
            d = _describe_image_part(img, i)
            if d:
                yield d
    except Exception as e:
        print(f"  ~ OCR страницы PDF {i} не удался: {e}")


def _load_pdf(path: Path):
    import fitz  # pymupdf
    doc = fitz.open(path)
    # OCR страниц без текстового слоя (сканы/дизайн-страницы) — по флагу OCR_IMAGES
    ocr_on = _enabled("OCR_IMAGES")
    try:
        min_chars = int(_ocr_setting("OCR_MIN_CHARS", 25))
    except Exception:
        min_chars = 25
    try:
        for i, page in enumerate(doc, 1):
            txt = page.get_text("text")
            if txt.strip():
                yield {"text": txt, "page": i}
            # мало или нет текста — вероятно, страница нарисована картинкой: распознаём
            if ocr_on and len(txt.strip()) < min_chars and _ocr_available():
                yield from _ocr_pdf_page(page, i)
    finally:
        # закрываем документ в finally (в т. ч. при закрытии генератора до конца) —
        # иначе утекают файловые дескрипторы PyMuPDF при массовой индексации.
        try:
            doc.close()
        except Exception:
            pass


def _load_docx(path: Path):
    import docx
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    # таблицы внутри документа
    for tbl in d.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    if parts:
        yield {"text": "\n".join(parts), "page": None}


def _load_pptx(path: Path):
    from pptx import Presentation
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, 1):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text)
        # заметки докладчика часто = расшифровка обучения
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            chunks.append("Заметки: " + slide.notes_slide.notes_text_frame.text)
        if chunks:
            yield {"text": "\n".join(chunks), "page": i}


def _read_text_any(path: Path) -> str:
    """Прочитать текстовый файл, определив кодировку.
    Порядок: charset_normalizer → utf-8-sig → cp1251 → latin-1 → utf-8(replace).
    Общий хелпер для ВСЕХ текстовых загрузчиков: read_text(errors='ignore') молча
    ломает кириллицу в cp1251, поэтому кодировку надо подбирать, а не игнорировать."""
    raw = Path(path).read_bytes()
    encodings: list[str] = []
    try:
        import charset_normalizer
        best = charset_normalizer.from_bytes(raw).best()
        if best and best.encoding:
            encodings.append(best.encoding)
    except Exception:
        pass
    for e in ("utf-8-sig", "cp1251", "latin-1"):
        if e not in encodings:
            encodings.append(e)
    for e in encodings:
        try:
            return raw.decode(e)
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


def _read_csv_any(path: Path, pd):
    """Прочитать CSV с авто-определением кодировки и разделителя.
    Поддерживает UTF-8/UTF-8-BOM, Windows-1251 (кириллица), Latin-1 и др.;
    разделитель — , ; \\t |."""
    import io
    # 1) кодировка: общий хелпер (charset_normalizer + cp1251-фолбэк)
    text = _read_text_any(path)
    # 2) разделитель: Sniffer, иначе — самый частый в первой строке
    import csv as _csv
    sample = text[:8192]
    sep = ","
    try:
        sep = _csv.Sniffer().sniff(sample, delimiters=";,\t|").delimiter
    except Exception:
        first = (sample.splitlines() or [""])[0]
        cand = max((";", ",", "\t", "|"), key=lambda d: first.count(d))
        if first.count(cand) > 0:
            sep = cand
    return pd.read_csv(io.StringIO(text), dtype=str, keep_default_na=False,
                       sep=sep, engine="python")


def _load_table(path: Path):
    """Прайс-листы и таблицы: каждую строку превращаем в 'колонка: значение'."""
    import pandas as pd
    if path.suffix.lower() == ".csv":
        frames = {"csv": _read_csv_any(path, pd)}
    else:
        frames = pd.read_excel(path, sheet_name=None, dtype=str)
    for sheet, df in frames.items():
        df = df.fillna("")
        rows = []
        for _, row in df.iterrows():
            pairs = [f"{col}: {val}" for col, val in row.items() if str(val).strip()]
            if pairs:
                rows.append("; ".join(pairs))
        if rows:
            yield {"text": f"Лист «{sheet}»\n" + "\n".join(rows), "page": None}


def _load_html(path: Path):
    from bs4 import BeautifulSoup
    ext = path.suffix.lower()
    if ext in {".mhtml", ".mht"}:
        # MHTML — это MIME-архив; вытаскиваем html-часть
        import email
        msg = email.message_from_bytes(path.read_bytes())
        html = ""
        for part in msg.walk():
            if part.get_content_type() == "text/html":
                payload = part.get_payload(decode=True) or b""
                charset = part.get_content_charset() or "utf-8"
                html += payload.decode(charset, errors="ignore")
        raw = html or _read_text_any(path)
    else:
        raw = _read_text_any(path)
    soup = BeautifulSoup(raw, "html.parser")
    text = soup.get_text(separator="\n")
    if text.strip():
        yield {"text": text, "page": None}


def find_oda_converter():
    """Найти ODA File Converter: путь из настроек/env, затем PATH, затем типовые места
    установки (в т. ч. .app на macOS). Возвращает путь к исполняемому файлу или None."""
    import os as _os
    import shutil
    cands = []
    try:
        import settings
        p = (settings.get("ODA_CONVERTER_PATH") or "").strip()
    except Exception:
        p = _os.getenv("ODA_CONVERTER_PATH", "")
    if p:
        cands.append(p)
    w = shutil.which("ODAFileConverter")
    if w:
        cands.append(w)
    cands += [
        "/Applications/ODAFileConverter.app/Contents/MacOS/ODAFileConverter",  # macOS
        "/usr/bin/ODAFileConverter", "/usr/local/bin/ODAFileConverter",
        "/opt/ODAFileConverter/ODAFileConverter",
    ]
    for c in cands:
        try:
            if c and Path(c).exists():
                return c
        except Exception:
            continue
    return None


def _dwg_to_dxf(path: Path):
    """Конвертировать DWG -> DXF. Сначала dwg2dxf (libredwg), при неудаче —
    ODA File Converter (если установлен). Возвращает путь к DXF или None
    (не бросает исключение — вызывающий перейдёт к аварийному извлечению)."""
    import os as _os
    import shutil
    import subprocess
    import tempfile
    # Уникальный временный DXF на каждую конверсию: при INGEST_WORKERS>1 два потока с
    # одинаковым stem писали бы в один и тот же файл (gettempdir()/stem+"_conv.dxf") —
    # гонка и порча данных. mkstemp даёт уникальное имя; вызывающий удалит его (tmp.unlink).
    _fd, _out = tempfile.mkstemp(suffix=".dxf", prefix=path.stem + "_conv_")
    _os.close(_fd)
    out = Path(_out)
    if shutil.which("dwg2dxf"):
        try:
            subprocess.run(["dwg2dxf", "-o", str(out), str(path)],
                           check=True, capture_output=True, timeout=180)
        except Exception:
            pass                       # частый случай: неподдерживаемая версия DWG
        if out.exists() and out.stat().st_size > 0:
            return out
    # запасной конвертер — ODA File Converter (директория→директория)
    oda = find_oda_converter()
    if oda:
        ind = outd = None
        try:
            ind, outd = Path(tempfile.mkdtemp()), Path(tempfile.mkdtemp())
            shutil.copy(path, ind / path.name)
            cmd = [oda, str(ind), str(outd), "ACAD2018", "DXF", "0", "1", "*.DWG"]
            # на Linux без дисплея ODA (Qt) требует виртуальный X-сервер
            if _os.name == "posix" and not _os.environ.get("DISPLAY") \
                    and shutil.which("xvfb-run") and "darwin" not in _os.sys.platform:
                cmd = ["xvfb-run", "-a"] + cmd
            subprocess.run(cmd, capture_output=True, timeout=300)
            cand = outd / (path.stem + ".dxf")
            if cand.exists() and cand.stat().st_size > 0:
                shutil.copy(cand, out)
                return out
        except Exception:
            pass
        finally:                            # не оставляем временные папки ODA (утечка при массе DWG)
            for _d in (ind, outd):
                if _d is not None:
                    shutil.rmtree(_d, ignore_errors=True)
    # конверсия не удалась — убираем пустой временный файл, чтобы не копить мусор
    try:
        out.unlink()
    except Exception:
        pass
    return None


def _dwg_scrape_text(path) -> str:
    """Аварийное извлечение текста прямо из бинарника DWG, когда конвертер (dwg2dxf/ODA)
    не справился. Достаёт печатаемые строки (cp1251 и UTF-16LE) и оставляет только
    осмысленные: со словами/пробелами/кириллицей или артикулы (буквы+цифры, напр.
    SPL-3-A-105). Грубо, но лучше, чем потерять файл целиком."""
    import re
    try:
        data = Path(path).read_bytes()
    except Exception:
        return ""
    chunks = []
    for enc in ("cp1251", "utf-16-le"):
        try:
            s = data.decode(enc, errors="ignore")
        except Exception:
            continue
        # режем на «слова» по непечатаемым символам
        chunks += re.split(r"[\x00-\x1f\x7f]+", s)
    out = []
    for v in chunks:
        v = v.strip()
        if len(v) < 3:
            continue
        looks_partno = bool(re.search(r"\d", v)) and bool(re.search(r"[A-Za-zА-Яа-я]", v)) \
            and ("-" in v or "_" in v)
        if _has_cyr_or_space(v) or looks_partno:
            # отсекаем очевидный техномусор классов DWG
            if not re.match(r"^(Ac[A-Z]|AutoCAD|ANSI|ISO-|\*|SHX|Standard$)", v):
                out.append(v)
        if sum(len(x) for x in out) > 8000:
            break
    return "\n".join(dict.fromkeys(out)).strip()


def _dxf_scrape_text(src) -> str:
    """Аварийное извлечение текста из повреждённого DXF без полного парсинга.
    DXF — текстовый формат «код/значение»; текстовые строки идут под группами 1 и 3
    (TEXT/MTEXT/ATTRIB). Читаем их напрямую — работает даже когда ezdxf падает
    («Invalid group code», «Invalid transformation matrix» и т. п.)."""
    import re
    try:
        raw = Path(src).read_text(encoding="utf-8", errors="ignore")
    except Exception:
        try:
            raw = Path(src).read_bytes().decode("cp1251", errors="ignore")
        except Exception:
            return ""
    lines = raw.splitlines()
    vals = []
    for i in range(len(lines) - 1):
        if lines[i].strip() in ("1", "3"):        # 1 — основной текст, 3 — доп. MTEXT
            v = lines[i + 1].strip()
            if v and any(ch.isalpha() for ch in v):
                vals.append(v)
    cleaned = []
    for v in vals:
        v = re.sub(r"\\[A-Za-z][^;\\]*;", "", v)   # форматные коды MTEXT: \A1; \fArial|...;
        v = v.replace("\\P", " ").replace("\\~", " ")
        v = re.sub(r"[{}]", "", v).strip()
        if not v or len(v) < 2:
            continue
        # берём: текст со словами/пробелами/кириллицей ИЛИ артикулы (буквы+цифры, напр. SPL-3-D-60);
        # отбрасываем технический мусор (имена стилей/типов линий: Standard, ByLayer, txt.shx …)
        looks_partno = bool(re.search(r"\d", v)) and bool(re.search(r"[A-Za-zА-Яа-я]", v)) \
            and ("-" in v or "_" in v)
        if _has_cyr_or_space(v) or looks_partno:
            cleaned.append(v)
    return "\n".join(dict.fromkeys(cleaned)).strip()


def _has_cyr_or_space(s: str) -> bool:
    return (" " in s) or any("а" <= ch.lower() <= "я" or ch.lower() == "ё" for ch in s)


_CAD_VISION_PROMPT = (
    "Это инженерный/технический чертёж (САПР). Опиши по-русски, что на нём изображено: "
    "тип изделия/детали/узла, основные элементы и их назначение, характерные размеры и "
    "обозначения, материалы, надписи и содержимое основной надписи (штампа), если видны. "
    "Кратко и по делу, без вступлений.")


def _render_cad_image(doc):
    """Отрисовать модель чертежа (DXF/сконвертированный DWG) в растровое изображение
    (PIL Image) для описания vision-моделью. None при неудаче/пустом чертеже."""
    try:
        import matplotlib
        matplotlib.use("Agg")                       # без дисплея (сервер)
        import matplotlib.pyplot as plt
        from ezdxf.addons.drawing import RenderContext, Frontend
        from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
        from PIL import Image
        import io as _io
        msp = doc.modelspace()
        fig = plt.figure(figsize=(14, 10))
        ax = fig.add_axes([0, 0, 1, 1]); ax.set_axis_off()
        try:
            Frontend(RenderContext(doc), MatplotlibBackend(ax)).draw_layout(msp, finalize=True)
        except Exception:
            plt.close(fig); return None
        buf = _io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return Image.open(buf).convert("RGB")
    except Exception as e:
        print(f"  ~ CAD: рендер чертежа в изображение не удался: {e}")
        return None


# Сентинел начала секции «Image Data» (превью) в DWG (R13…R2018), 16 байт.
_DWG_IMG_SENTINEL = bytes(
    (0x1F, 0x25, 0x6D, 0x07, 0xD4, 0xDF, 0x43, 0xB1,
     0x8A, 0x3F, 0x4C, 0xDE, 0x35, 0x5C, 0x65, 0xEB))


def _dwg_preview_image(path):
    """Достать встроенное превью-изображение из DWG (когда сам чертёж не
    конвертируется). Современный AutoCAD хранит превью как PNG, старый — как BMP
    (DIB без файлового заголовка). Возвращает PIL.Image (RGB) или None.

    Сначала парсим секцию Image Data по сентинелу, при неудаче — грубо ищем
    сигнатуру PNG в файле (версионно-независимый фолбэк)."""
    import struct
    try:
        from PIL import Image
        import io as _io
    except Exception:
        return None
    try:
        data = Path(path).read_bytes()
    except Exception:
        return None

    def _open(buf):
        try:
            im = Image.open(_io.BytesIO(buf))
            im.load()
            if im.width < 8 or im.height < 8:      # заглушка/битое превью
                return None
            return im.convert("RGB")
        except Exception:
            return None

    def _from_entry(code, start, size):
        if size <= 0 or start <= 0 or start + size > len(data):
            return None
        blob = data[start:start + size]
        if code == 6:                              # PNG — готовый файл
            return _open(blob)
        if code == 2:                              # BMP как DIB — добавляем файл-заголовок
            try:
                hsize = struct.unpack_from("<I", blob, 0)[0]
                bits = struct.unpack_from("<H", blob, 14)[0]
                clr = struct.unpack_from("<I", blob, 32)[0]
                if clr == 0 and bits <= 8:
                    clr = 1 << bits
                px = 14 + hsize + clr * 4
                fh = b"BM" + struct.pack("<IHHI", 14 + size, 0, 0, px)
                return _open(fh + blob)
            except Exception:
                return None
        return None

    def _parse_section(p):
        """Разобрать секцию Image Data начиная со смещения p (на сентинел)."""
        try:
            # sentinel — 16 байт; проверяем только сигнатуру начала (1F 25 6D 07 D4)
            if data[p:p + 5] != _DWG_IMG_SENTINEL[:5]:
                return None
            q = p + 16 + 4                         # +сентинел +overall size(RL)
            count = data[q]; q += 1
            if not (0 < count < 16):
                return None
            for _ in range(count):
                code = data[q]; q += 1
                start, size = struct.unpack_from("<II", data, q); q += 8
                im = _from_entry(code, start, size)
                if im is not None:
                    return im
        except Exception:
            return None
        return None

    # --- способ 1: адрес превью из заголовка DWG (RL @ 0x0D), R2000+ ---
    try:
        addr = struct.unpack_from("<I", data, 0x0D)[0]
        if 0 < addr < len(data) - 32:
            im = _parse_section(addr)
            if im is not None:
                return im
    except Exception:
        pass

    # --- способ 2: поиск сентинела секции Image Data по файлу ---
    try:
        p = data.find(_DWG_IMG_SENTINEL[:5])
        while p >= 0:
            im = _parse_section(p)
            if im is not None:
                return im
            p = data.find(_DWG_IMG_SENTINEL[:5], p + 1)
    except Exception:
        pass

    # --- способ 3 (фолбэк): найти PNG-сигнатуру и вырезать до IEND ---
    try:
        sig = b"\x89PNG\r\n\x1a\n"
        i = data.find(sig)
        if i >= 0:
            j = data.find(b"IEND", i)
            if j >= 0:
                im = _open(data[i:j + 8])          # IEND + 4 байта CRC
                if im is not None:
                    return im
    except Exception:
        pass
    return None


def _describe_dwg_preview(path):
    """Описать неконвертируемый DWG по встроенному превью через vision-модель.
    → {'text','vision_desc'} или None."""
    if not _vision_available():
        return None
    img = _dwg_preview_image(path)
    if img is None:
        print(f"  ~ CAD: {path.name} — встроенное превью в DWG не найдено, "
              f"описать изображением нечем")
        return None
    try:
        import llm_backend
        desc = llm_backend.describe_image(img, prompt=_CAD_VISION_PROMPT)
        _vision_result(bool(desc and desc.strip()))
        if desc and desc.strip():
            print(f"  ~ CAD: {path.name} — DWG описан vision-моделью по встроенному превью")
            return {"text": "Описание чертежа DWG (vision-модель, встроенное превью):\n"
                    + desc.strip(), "page": None, "vision_desc": True}
    except Exception as e:
        _vision_result(False)
        print(f"  ~ CAD: {path.name} — описание превью DWG моделью не удалось: {e}")
    return None


def _describe_cad_part(doc):
    """Отрисовать чертёж и описать vision-моделью → {'text','vision_desc'} или None."""
    if not _vision_available():
        return None
    img = _render_cad_image(doc)
    if img is None:
        return None
    try:
        import llm_backend
        desc = llm_backend.describe_image(img, prompt=_CAD_VISION_PROMPT)
        _vision_result(bool(desc and desc.strip()))
        if desc and desc.strip():
            return {"text": "Описание чертежа (vision-модель):\n" + desc.strip(),
                    "page": None, "vision_desc": True}
    except Exception as e:
        _vision_result(False)
        print(f"  ~ CAD: описание чертежа моделью не удалось: {e}")
    return None


def _load_cad(path: Path):
    """Извлечь весь текст из чертежа DXF/DWG: TEXT, MTEXT, атрибуты блоков,
    размеры, имена слоёв (для DWG требуется конвертация в DXF). Повреждённые файлы
    (частый случай после конвертации DWG→DXF) обрабатываются по возможности, с
    аварийным текстовым извлечением как фолбэком."""
    import ezdxf
    # приглушаем шумные предупреждения ezdxf (Found non-unique entity handle,
    # Missing required ENDBLK, Invalid group code …) — они не мешают извлечению текста
    import logging as _lg
    _lg.getLogger("ezdxf").setLevel(_lg.ERROR)
    src, tmp = path, None
    if path.suffix.lower() == ".dwg":
        src = _dwg_to_dxf(path)
        tmp = src
        if src is None:
            # конвертер не справился (частый случай для новых версий DWG) —
            # не роняем файл ошибкой, а пробуем достать текст из бинарника
            scraped = _dwg_scrape_text(path)
            if scraped:
                print(f"  ~ CAD: {path.name} — DWG не конвертируется, "
                      f"аварийное извлечение текста из бинарника")
                yield {"text": "Текст чертежа DWG (аварийное извлечение):\n" + scraped,
                       "page": None}
            elif not _enabled("CAD_DWG_DESCRIBE"):
                print(f"  ~ CAD: {path.name} — DWG не конвертируется и текст не найден "
                      f"(нужен libredwg/ODA или сохраните как DXF/PDF); пропуск")
            # Неконвертируемый DWG → описать по встроенному превью vision-моделью
            # (опция CAD_DWG_DESCRIBE). Работает и когда текст не извлёкся.
            if _enabled("CAD_DWG_DESCRIBE"):
                d = _describe_dwg_preview(path)
                if d:
                    yield d
            return
    try:
        doc = None
        try:
            doc = ezdxf.readfile(str(src))
        except Exception:
            try:
                from ezdxf import recover
                doc, _ = recover.readfile(str(src))
            except Exception as e:
                print(f"  ~ CAD: {path.name} — ezdxf не смог разобрать ({e}); "
                      f"аварийное извлечение текста")
                doc = None

        lines = []

        def grab(container):
            it = iter(container)
            while True:
                try:
                    e = next(it)
                except StopIteration:
                    break
                except Exception:
                    break            # итератор повреждён — прекращаем этот контейнер
                try:
                    t = e.dxftype()
                    s = ""
                    if t in ("TEXT", "ATTRIB", "ATTDEF"):
                        s = e.dxf.get("text", "")
                    elif t == "MTEXT":
                        s = e.text
                    elif t == "DIMENSION":
                        s = e.dxf.get("text", "")
                        if s in ("", "<>"):
                            s = ""
                    if s and s.strip():
                        lines.append(s.strip())
                    if t == "INSERT":  # атрибуты вставленных блоков
                        for a in getattr(e, "attribs", []) or []:
                            v = a.dxf.get("text", "")
                            if v and v.strip():
                                lines.append(v.strip())
                except Exception:
                    continue

        layers = []
        if doc is not None:
            try:
                for layout in doc.layouts:      # модель + листы
                    try:
                        grab(layout)
                    except Exception:
                        continue
            except Exception:
                pass
            try:
                for blk in doc.blocks:          # текст внутри определений блоков
                    try:
                        grab(blk)
                    except Exception:
                        continue
            except Exception:
                pass
            try:
                layers = [ly.dxf.name for ly in doc.layers]
            except Exception:
                layers = []

        body = "\n".join(dict.fromkeys(lines)).strip()  # дедуп с сохранением порядка
        # фолбэк: ezdxf не дал текста (или упал) — скребём DXF как текст
        if not body:
            scraped = _dxf_scrape_text(src)
            if scraped:
                body = "Текст чертежа (аварийное извлечение):\n" + scraped
        if layers:
            body += ("\n" if body else "") + "Слои: " + ", ".join(layers[:300])
        body = body.strip()
        if body:
            yield {"text": body, "page": None}
        # Дополнительно: отрисовать чертёж и описать vision-моделью (опция CAD_LLM_DESCRIBE).
        # Делаем ДО очистки временного DXF (в finally) — нужен разобранный doc.
        if doc is not None and _enabled("CAD_LLM_DESCRIBE"):
            d = _describe_cad_part(doc)
            if d:
                yield d
    finally:
        if tmp:
            try:
                tmp.unlink()
            except Exception:
                pass


def _load_cad_exchange(path: Path):
    """Текстовые метаданные из STEP (.stp/.step) и IGES (.igs/.iges):
    названия деталей/изделий, описания, заголовок, единицы, автор.
    Геометрия не извлекается (для текстового поиска она бесполезна)."""
    ext = path.suffix.lower()
    text = _read_text_any(path)
    body = _parse_iges(text) if ext in (".igs", ".iges") else _parse_step(text)
    if body.strip():
        yield {"text": body, "page": None}


def _parse_step(t: str) -> str:
    import re
    seen, names = set(), []
    # все строковые литералы STEP в одинарных кавычках ('' = апостроф)
    for m in re.finditer(r"'((?:[^']|'')*)'", t):
        s = m.group(1).replace("''", "'").strip()
        if len(s) >= 2 and not s.isdigit() and re.search(r"[A-Za-zА-Яа-я0-9]", s):
            if s not in seen:
                seen.add(s)
                names.append(s)
    if not names:
        return ""
    return "Метаданные STEP (названия, описания, единицы, заголовок):\n" + "\n".join(names[:2000])


def _parse_iges(t: str) -> str:
    import re
    start, glob = [], []
    for line in t.splitlines():
        if len(line) >= 73:
            sec = line[72]
            if sec == "S":  # Start section — свободный текст-описание
                s = line[:72].strip()
                if s:
                    start.append(s)
            elif sec == "G":  # Global section — параметры (Hollerith-строки)
                glob.append(line[:72])
    gtext = "".join(glob)
    holler, i = [], 0
    while i < len(gtext):  # Hollerith: NNNNH<строка ровно NNNN символов>
        m = re.match(r"(\d+)H", gtext[i:])
        if m:
            n = int(m.group(1))
            pos = i + m.end()
            val = gtext[pos:pos + n].strip()
            if val and re.search(r"[A-Za-zА-Яа-я0-9]", val):
                holler.append(val)
            i = pos + n
        else:
            i += 1
    out = []
    if start:
        out.append("Описание (IGES Start):\n" + "\n".join(start))
    if holler:
        out.append("Параметры (IGES Global):\n" + ", ".join(holler))
    return "\n".join(out)


_OCR_OK = None  # кеш проверки доступности OCR


def _ocr_available() -> bool:
    """Установлены ли pytesseract + сам tesseract. Предупреждаем один раз."""
    global _OCR_OK
    if _OCR_OK is None:
        import importlib.util
        import shutil
        has_lib = importlib.util.find_spec("pytesseract") is not None
        has_bin = shutil.which("tesseract") is not None
        _OCR_OK = has_lib and has_bin
        if not _OCR_OK:
            miss = []
            if not has_lib:
                miss.append("pytesseract (pip install pytesseract Pillow)")
            if not has_bin:
                miss.append("tesseract (системный пакет, напр. tesseract-ocr + -rus)")
            print(f"  ~ OCR недоступен — пропускаю распознавание картинок/RAW. "
                  f"Не хватает: {', '.join(miss)}. Либо отключите OCR в настройках.")
    return _OCR_OK


def _ocr_setting(key, default):
    try:
        v = settings.get(key)
        return default if v in (None, "") else v
    except Exception:
        return default


def _ocr_lang():
    """Языки OCR: явная настройка OCR_LANGS, иначе автоопределение."""
    explicit = _ocr_setting("OCR_LANGS", "")
    if explicit:
        return str(explicit)
    import pytesseract
    try:
        langs = pytesseract.get_languages(config="")
        if "rus" in langs:
            return "rus+eng" if "eng" in langs else "rus"
    except Exception:
        pass
    return "eng"


def _ocr_config():
    """Строка конфигурации tesseract из настроек OEM/PSM."""
    try:
        oem = int(_ocr_setting("OCR_OEM", 3))
        psm = int(_ocr_setting("OCR_PSM", 3))
        return f"--oem {oem} --psm {psm}"
    except Exception:
        return ""


def _ocr_image(img):
    """PIL.Image → OCR → текстовый (searchable) PDF → извлечённый текст.
    Общий помощник для растровых изображений и RAW-фото. Параметры распознавания
    (язык, OEM/PSM, макс. размер, предобработка) берутся из настроек OCR."""
    import io
    import pytesseract
    import fitz  # pymupdf

    # уменьшаем огромные снимки — ускоряет OCR без потери читаемости текста
    try:
        maxdim = int(_ocr_setting("OCR_MAX_DIM", 3500))
    except Exception:
        maxdim = 3500
    if maxdim > 0 and max(img.size) > maxdim:
        k = maxdim / max(img.size)
        img = img.resize((max(1, int(img.size[0] * k)), max(1, int(img.size[1] * k))))
    # предобработка: оттенки серого + автоконтраст (помогает на шумных сканах)
    if _ocr_setting("OCR_PREPROCESS", False):
        try:
            from PIL import ImageOps
            img = ImageOps.autocontrast(img.convert("L"))
        except Exception:
            pass
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")

    pdf_bytes = pytesseract.image_to_pdf_or_hocr(
        img, lang=_ocr_lang(), extension="pdf", config=_ocr_config())
    doc = fitz.open(stream=io.BytesIO(pdf_bytes).getvalue(), filetype="pdf")
    try:
        for i, page in enumerate(doc, 1):
            txt = page.get_text("text")
            if txt.strip():
                yield {"text": txt, "page": i}
    finally:
        try:
            doc.close()   # не утекают дескрипторы PyMuPDF
        except Exception:
            pass


def _load_image(path: Path):
    """Растровое изображение (jpg/png/…) → OCR. Полезно для сканов и фото
    документов, скриншотов прайсов и т. п. Если OCR дал ≤1 чанк — изображение
    дополнительно описывается vision-моделью (если включено)."""
    describe_on = _enabled("OCR_LLM_DESCRIBE")
    ocr_on = _ocr_available()
    if not ocr_on and not describe_on:
        return  # ни OCR, ни описания — не декодируем картинку зря
    from PIL import Image

    if ocr_on:
        print(f"  ~ распознаю (OCR) {path.name} ...")
    with Image.open(path) as img:
        img.load()
        total = 0
        if ocr_on:
            for part in _ocr_image(img):
                t = (part.get("text") or "").strip()
                if t:
                    total += len(t)
                yield part
        if _describe_due(total):
            print(f"  ~ описываю изображение моделью: {path.name}")
            d = _describe_image_part(img)
            if d:
                yield d


def _load_raw(path: Path):
    """RAW-фото (CR2 и др.) → изображение → OCR → текст.
    Полезно для сфотографированных документов. При ≤1 чанке — описание моделью."""
    describe_on = _enabled("OCR_LLM_DESCRIBE")
    ocr_on = _ocr_available()
    if not ocr_on and not describe_on:
        return
    import rawpy
    from PIL import Image

    if ocr_on:
        print(f"  ~ распознаю (OCR) {path.name} ...")
    with rawpy.imread(str(path)) as raw:
        rgb = raw.postprocess()
    img = Image.fromarray(rgb)
    total = 0
    if ocr_on:
        for part in _ocr_image(img):
            t = (part.get("text") or "").strip()
            if t:
                total += len(t)
            yield part
    if _describe_due(total):
        print(f"  ~ описываю изображение моделью: {path.name}")
        d = _describe_image_part(img)
        if d:
            yield d


def _load_svg(path: Path):
    """SVG — извлекаем текст из элементов <text>/<tspan>."""
    import re
    data = _read_text_any(path)
    # вытаскиваем содержимое текстовых тегов
    parts = re.findall(r"<(?:text|tspan)\b[^>]*>(.*?)</(?:text|tspan)>", data,
                       flags=re.DOTALL | re.IGNORECASE)
    text = "\n".join(re.sub(r"<[^>]+>", " ", p) for p in parts).strip()
    if text:
        yield {"text": text, "page": None}


def _load_xml(path: Path):
    """XML — собираем весь видимый текст из узлов."""
    try:
        import xml.etree.ElementTree as ET
        root = ET.parse(str(path)).getroot()
        parts = [t.strip() for t in root.itertext() if t and t.strip()]
        text = "\n".join(parts)
    except Exception:
        import re
        raw = _read_text_any(path)
        text = re.sub(r"<[^>]+>", " ", raw)
    if text.strip():
        yield {"text": text, "page": None}


def _load_json(path: Path):
    """JSON — плоское текстовое представление пар ключ/значение."""
    import json
    data = json.loads(_read_text_any(path))

    def walk(obj, prefix=""):
        if isinstance(obj, dict):
            for k, v in obj.items():
                yield from walk(v, f"{prefix}{k}: ")
        elif isinstance(obj, list):
            for it in obj:
                yield from walk(it, prefix)
        else:
            s = str(obj).strip()
            if s:
                yield f"{prefix}{s}"

    text = "\n".join(walk(data))
    if text.strip():
        yield {"text": text, "page": None}


def _load_url(path: Path):
    """Ярлык .url (Windows Internet Shortcut) — извлекаем адрес ссылки."""
    import re
    data = _read_text_any(path)
    m = re.search(r"URL\s*=\s*(\S+)", data, flags=re.IGNORECASE)
    url = m.group(1).strip() if m else ""
    text = f"Ссылка ({path.stem}): {url}".strip()
    if url:
        yield {"text": text, "page": None}


def _load_doc(path: Path):
    """Старый Word (.doc) — конвертация через antiword или LibreOffice."""
    import shutil
    import subprocess
    import tempfile

    # 1) antiword — быстрый и точный для .doc
    if shutil.which("antiword"):
        try:
            out = subprocess.run(["antiword", str(path)], capture_output=True,
                                 timeout=120)
            text = out.stdout.decode("utf-8", errors="ignore")
            if text.strip():
                yield {"text": text, "page": None}
                return
        except Exception:
            pass

    # 2) LibreOffice/soffice — конвертируем .doc → .docx, читаем как docx
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        with tempfile.TemporaryDirectory() as td:
            try:
                subprocess.run([soffice, "--headless", "--convert-to", "docx",
                                "--outdir", td, str(path)],
                               capture_output=True, timeout=240)
                conv = Path(td) / (path.stem + ".docx")
                if conv.exists():
                    yield from _load_docx(conv)
                    return
            except Exception:
                pass
    print(f"  ! пропуск {path.name}: нет antiword/LibreOffice для .doc")


def _load_msg(path: Path):
    """Письмо Outlook (.msg) — тема, отправитель, тело."""
    import extract_msg
    m = extract_msg.Message(str(path))
    try:
        head = []
        for label, val in (("Тема", m.subject), ("От", m.sender),
                           ("Кому", m.to), ("Дата", m.date)):
            if val:
                head.append(f"{label}: {val}")
        body = (m.body or "").strip()
        text = ("\n".join(head) + "\n\n" + body).strip()
    finally:
        try:
            m.close()
        except Exception:
            pass
    if text:
        yield {"text": text, "page": None}


class _ArchiveAbort(Exception):
    """Распаковка прервана из соображений безопасности (превышен лимит объёма или
    небезопасный путь члена). В отличие от прочих ошибок НЕ приводит к фолбэку на
    системные утилиты — иначе защита обходилась бы через `7z x`/`bsdtar`."""


def _safe_member_path(dest: Path, name: str) -> Path:
    """Целевой путь члена архива внутри dest. Отказ при абсолютных путях и '..'
    (path traversal). Возвращает нормализованный путь под dest."""
    droot = dest.resolve()
    target = (dest / name).resolve()
    if target != droot and droot not in target.parents:
        raise _ArchiveAbort(f"небезопасный путь в архиве: {name!r}")
    return target


def _copy_capped(src, dst, budget: list) -> None:
    """Потоковое копирование src→dst с учётом остатка бюджета байт (budget[0]).
    При превышении — _ArchiveAbort (прерывание распаковки без фолбэка)."""
    while True:
        block = src.read(1 << 20)
        if not block:
            break
        budget[0] -= len(block)
        if budget[0] < 0:
            raise _ArchiveAbort("распакованный объём превышает лимит")
        dst.write(block)


def _dir_size(dest: Path) -> int:
    total = 0
    for p in Path(dest).rglob("*"):
        try:
            if p.is_file():
                total += p.stat().st_size
        except Exception:
            pass
    return total


def _extract_archive(path: Path, dest: Path) -> bool:
    """Распаковать архив в каталог dest. Сначала пробуем библиотеки Python,
    затем системные утилиты (7z/bsdtar/unar). Возвращает True при успехе."""
    import shutil
    import subprocess
    ext = path.suffix.lower()

    # 1) Python-библиотеки
    try:
        if ext == ".zip":
            import zipfile
            budget = [_ARCHIVE_MAX_BYTES]
            with zipfile.ZipFile(path) as z:
                total = sum(i.file_size for i in z.infolist())
                if total > _ARCHIVE_MAX_BYTES:
                    raise _ArchiveAbort("распакованный объём превышает лимит")
                for info in z.infolist():
                    if info.is_dir():
                        _safe_member_path(dest, info.filename).mkdir(parents=True, exist_ok=True)
                        continue
                    target = _safe_member_path(dest, info.filename)
                    target.parent.mkdir(parents=True, exist_ok=True)
                    with z.open(info) as src, open(target, "wb") as dst:
                        _copy_capped(src, dst, budget)
            return True
        if ext == ".7z":
            import py7zr
            with py7zr.SevenZipFile(path, "r") as z:
                # py7zr извлекает потоково на диск, но не даёт удобного пофайлового
                # счётчика — оцениваем распакованный объём заранее по заголовкам.
                try:
                    total = sum((getattr(fi, "uncompressed", 0) or 0) for fi in z.list())
                except Exception:
                    total = 0
                if total > _ARCHIVE_MAX_BYTES:
                    raise _ArchiveAbort("распакованный объём превышает лимит")
                z.extractall(dest)
            # backstop: фактический объём на диске (на случай неверных заголовков)
            # FIXME(review): для 7z нет истинного потокового счётчика — контроль по
            # заголовкам + пост-проверка размера каталога, бомба может кратко лечь на диск.
            if _dir_size(dest) > _ARCHIVE_MAX_BYTES:
                raise _ArchiveAbort("распакованный объём превышает лимит")
            return True
        if ext == ".rar":
            import rarfile
            # rarfile сам по себе не распаковывает — ему нужен внешний бэкенд.
            # Подскажем доступный (unar/unrar/bsdtar), если авто-детект не сработал.
            for _t in ("unar", "unrar", "bsdtar"):
                if shutil.which(_t):
                    try:
                        rarfile.UNRAR_TOOL = _t
                    except Exception:
                        pass
                    break
            budget = [_ARCHIVE_MAX_BYTES]
            with rarfile.RarFile(path) as r:
                for info in r.infolist():
                    if info.isdir():
                        _safe_member_path(dest, info.filename).mkdir(parents=True, exist_ok=True)
                        continue
                    target = _safe_member_path(dest, info.filename)
                    target.parent.mkdir(parents=True, exist_ok=True)
                    with r.open(info) as src, open(target, "wb") as dst:
                        _copy_capped(src, dst, budget)
            return True
        if ext in {".tar", ".tgz"} or (ext in {".gz", ".bz2"} and ".tar" in path.name.lower()):
            import tarfile
            budget = [_ARCHIVE_MAX_BYTES]
            with tarfile.open(path) as t:
                # Ручная валидация членов (отказ при абсолютных путях/'..') вместо
                # молчаливого фолбэка на небезопасный extractall на Python < 3.11.4.
                for m in t:
                    if m.isdir():
                        _safe_member_path(dest, m.name).mkdir(parents=True, exist_ok=True)
                        continue
                    if not m.isfile():
                        continue   # симлинки/устройства/hardlink — пропускаем (безопасность)
                    target = _safe_member_path(dest, m.name)
                    target.parent.mkdir(parents=True, exist_ok=True)
                    src = t.extractfile(m)
                    if src is None:
                        continue
                    with src, open(target, "wb") as dst:
                        _copy_capped(src, dst, budget)
            return True
        if ext in {".gz", ".bz2"}:
            # одиночный поток (не tar-архив) — распаковываем в один файл
            import bz2
            import gzip
            opener = gzip.open if ext == ".gz" else bz2.open
            out = Path(dest) / path.stem  # отбрасываем .gz/.bz2
            budget = [_ARCHIVE_MAX_BYTES]
            with opener(path, "rb") as src, open(out, "wb") as dst:
                _copy_capped(src, dst, budget)
            return True
    except _ArchiveAbort as e:
        # лимит объёма/небезопасный член — НЕ пробуем системные утилиты (иначе обойдём
        # защиту от «архивных бомб» и path traversal), помечаем архив как нераспакованный.
        print(f"  ! {path.name}: распаковка прервана ({e})")
        return False
    except ImportError:
        pass  # нужной библиотеки нет — пробуем системные утилиты
    except Exception as e:
        print(f"  ! {path.name}: ошибка распаковки ({e}); пробую системные утилиты")

    # 2) системные утилиты (best-effort). Для .rar сначала специализированные
    # распаковщики (unar/unrar), затем универсальные (7z/bsdtar).
    tools = []
    if ext == ".rar":
        tools += [
            ["unar", "-quiet", "-force-overwrite", "-output-directory", str(dest), str(path)],
            ["unrar", "x", "-y", "-o+", str(path), str(dest) + "/"],
        ]
    tools += [
        ["7z", "x", "-y", f"-o{dest}", str(path)],
        ["7za", "x", "-y", f"-o{dest}", str(path)],
        ["bsdtar", "-xf", str(path), "-C", str(dest)],
        ["unar", "-quiet", "-output-directory", str(dest), str(path)],
    ]
    for tool in tools:
        if shutil.which(tool[0]):
            try:
                r = subprocess.run(tool, capture_output=True, timeout=1800)
                # успех: код 0 И в каталоге появились файлы
                if r.returncode == 0 and any(Path(dest).rglob("*")):
                    return True
            except Exception:
                continue
    return False


def _load_archive(path: Path, depth: int):
    """Распаковать архив и проиндексировать содержимое как обычные файлы.
    Внутренние файлы помечаются их путём внутри архива для цитирования."""
    if depth >= _ARCHIVE_MAX_DEPTH:
        print(f"  ! пропуск вложенного архива {path.name}: слишком глубоко")
        return
    import tempfile
    print(f"  ~ распаковываю архив {path.name} ...")
    with tempfile.TemporaryDirectory() as td:
        dest = Path(td)
        if not _extract_archive(path, dest):
            import shutil as _sh
            _have = any(_sh.which(t) for t in ("unar", "unrar", "bsdtar", "7z", "7za", "7zr"))
            if _have:
                print(f"  ! {path.name}: архив не распаковался — вероятно повреждён, обрезан "
                      f"(неполная загрузка) или это том многотомного RAR (нужны все части); пропуск")
            else:
                print(f"  ! не удалось распаковать {path.name}: нет утилит распаковки — "
                      f"установите: sudo apt install -y unar libarchive-tools p7zip-full")
            return
        n = 0
        for inner in sorted(dest.rglob("*")):
            if not inner.is_file():
                continue
            n += 1
            if n > _ARCHIVE_MAX_FILES:
                print(f"  ! {path.name}: слишком много файлов в архиве, остановка")
                break
            rel = inner.relative_to(dest)
            try:
                for part in load_file(inner, _depth=depth + 1):
                    txt = part.get("text", "")
                    if txt.strip():
                        # помечаем источник внутри архива — пригодится для цитат
                        yield {"text": f"[{path.name} → {rel}]\n{txt}",
                               "page": part.get("page")}
            except Exception as e:
                print(f"  ! {path.name} → {rel}: {e}")


_FASTER_WHISPER = None  # ленивый кеш модели faster-whisper
_WHISPER_DEV = None     # на каком устройстве собрана кэш-модель ("cuda"/"cpu")


def _is_cuda_oom(e) -> bool:
    s = str(e).lower()
    return ("out of memory" in s or "cuda failed" in s or "cublas" in s
            or "cudnn" in s or "cuda error" in s)


def _empty_cuda() -> None:
    """Отпустить кэш VRAM (после OOM/смены устройства), чтобы не тормозить эмбеддер."""
    try:
        import gc
        import torch
        gc.collect()
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
    except Exception:
        pass


def _free_whisper() -> None:
    global _FASTER_WHISPER
    _FASTER_WHISPER = None
    _empty_cuda()


def _faster_whisper_model(model, device):
    from faster_whisper import WhisperModel
    compute = "float16" if device == "cuda" else "int8"
    return WhisperModel(model, device=device, compute_type=compute)


def _av_windows(segments):
    """Сгруппировать сегменты Whisper (start/end/text) в окна ~CHUNK_SIZE символов,
    сохраняя тайминги начала/конца окна — чтобы потом показать кадр/фрагмент видео."""
    try:
        win = int(settings.get("CHUNK_SIZE"))
    except Exception:
        win = 900
    buf, t0, t1, n = [], None, None, 0
    for seg in segments:
        st = float(seg.get("start") or 0.0)
        en = float(seg.get("end") or st)
        tx = (seg.get("text") or "").strip()
        if not tx:
            continue
        if t0 is None:
            t0 = st
        t1 = en
        buf.append(tx)
        n += len(tx) + 1
        if n >= win:
            yield {"text": " ".join(buf), "page": None, "t_start": t0, "t_end": t1}
            buf, t0, t1, n = [], None, None, 0
    if buf:
        yield {"text": " ".join(buf), "page": None, "t_start": t0, "t_end": t1}


def transcribe_audio(path) -> str:
    """Распознать речь из аудиофайла настроенным бэкендом Whisper и вернуть текст.
    Используется, например, Телеграм-ботом для голосовых сообщений."""
    parts = []
    try:
        for p in _load_av(Path(path)):
            t = (p.get("text") or "").strip()
            if t:
                parts.append(t)
    except Exception as e:
        print(f"  ! транскрибация {path}: {e}")
    return " ".join(parts).strip()


def _load_av(path: Path):
    """Транскрибация аудио/видео. Бэкенд зависит от настройки WHISPER_BACKEND:
       mlx    — Apple Metal (mlx-whisper),
       faster — GPU/CPU (faster-whisper, CTranslate2).
    Текст режется на окна с таймингами (t_start/t_end) для выдачи кадров/фрагментов."""
    print(f"  ~ транскрибирую {path.name} ...")
    device = settings.get("DEVICE")
    model = settings.get("WHISPER_MODEL")
    if settings.get("WHISPER_BACKEND") == "faster":
        global _FASTER_WHISPER, _WHISPER_DEV
        # Устройство: настройка WHISPER_DEVICE (auto=следовать DEVICE; cpu|cuda).
        wdev = (settings.get("WHISPER_DEVICE") or "auto").strip().lower()
        want = "cuda" if device == "cuda" else "cpu"
        if wdev in ("cpu", "cuda"):
            want = wdev
        # если уже вынужденно упали на CPU в этом прогоне — на GPU не возвращаемся
        if _WHISPER_DEV == "cpu":
            want = "cpu"
        # порядок попыток: желаемое устройство, затем CPU (фолбэк при CUDA OOM)
        tries = [want, "cpu"] if want == "cuda" else [want]
        last = None
        for dev in tries:
            try:
                if _FASTER_WHISPER is None or _WHISPER_DEV != dev:
                    _free_whisper()
                    _FASTER_WHISPER = _faster_whisper_model(model, dev)
                    _WHISPER_DEV = dev
                seg_iter, _info = _FASTER_WHISPER.transcribe(str(path))
                # материализуем здесь: транскрибация ленивая, ошибка (в т.ч. CUDA OOM)
                # возникает при итерации — ловим её ДО первого yield, чтобы фолбэк на CPU
                # не продублировал уже выданные окна
                segs = [{"start": s.start, "end": s.end, "text": s.text} for s in seg_iter]
                yield from _av_windows(segs)
                return
            except Exception as e:
                last = e
                if dev == "cuda" and _is_cuda_oom(e):
                    print("  ! Whisper: не хватает памяти GPU (её занял vLLM) — перехожу на "
                          "CPU до конца индексации. Ускорить: запускать ingest на свободной "
                          "GPU (CUDA_VISIBLE_DEVICES) или WHISPER_DEVICE=cpu.")
                    _free_whisper()
                    _WHISPER_DEV = "cpu"
                    continue
                raise
        if last:
            raise last
    else:  # mlx
        import mlx_whisper
        result = mlx_whisper.transcribe(str(path), path_or_hf_repo=model)
        segs = result.get("segments") or []
        if segs:
            yield from _av_windows(segs)
        elif result.get("text", "").strip():
            yield {"text": result["text"], "page": None}
